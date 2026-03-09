"""文件解析模块 - 支持多种文件格式解析为文本"""
import os
import io
import json
import tempfile
from typing import List, Dict, Any, Optional
from datetime import datetime

from core.logger import get_logger

logger = get_logger(__name__)


class FileParser:
    """文件解析器 - 将各种格式文件解析为文本内容"""

    # 支持的文件扩展名映射
    EXTENSION_PARSERS = {
        # 文档与文本
        '.txt': 'parse_text',
        '.md': 'parse_markdown',
        '.markdown': 'parse_markdown',
        '.pdf': 'parse_pdf',
        '.html': 'parse_html',
        '.htm': 'parse_html',
        '.docx': 'parse_docx',
        # 数据与表格
        '.xlsx': 'parse_excel',
        '.xls': 'parse_excel',
        '.csv': 'parse_csv',
        # 其他格式
        '.eml': 'parse_eml',
        '.msg': 'parse_msg',
        '.pptx': 'parse_pptx',
        '.ppt': 'parse_pptx',
        '.xml': 'parse_xml',
        '.epub': 'parse_epub',
        # 图片
        '.jpg': 'parse_image',
        '.jpeg': 'parse_image',
        '.png': 'parse_image',
        '.gif': 'parse_image',
    }

    @staticmethod
    def parse(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析文件并返回文本片段列表

        Args:
            file_content: 文件字节内容
            filename: 文件名

        Returns:
            文本片段列表，每项包含 content, metadata
        """
        ext = os.path.splitext(filename)[1].lower()
        parser_method = FileParser.EXTENSION_PARSERS.get(ext)

        if not parser_method:
            logger.warning(f"不支持的文件格式: {ext}")
            return [{"content": f"[不支持的格式: {ext}]", "metadata": {"source": filename, "type": ext}}]

        parser = getattr(FileParser, parser_method, None)
        if not parser:
            return [{"content": f"[解析方法不存在: {ext}]", "metadata": {"source": filename, "type": ext}}]

        try:
            return parser(file_content, filename)
        except Exception as e:
            logger.error(f"解析文件失败: {filename}, 错误: {e}")
            return [{"content": f"[解析失败: {str(e)}]", "metadata": {"source": filename, "type": ext, "error": str(e)}}]

    @staticmethod
    def parse_text(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析纯文本文件"""
        text = file_content.decode('utf-8', errors='ignore')
        return [{"content": text, "metadata": {"source": filename, "type": ".txt"}}]

    @staticmethod
    def parse_markdown(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 Markdown 文件"""
        import markdown
        text = file_content.decode('utf-8', errors='ignore')
        # 转换为 HTML 再提取文本
        html = markdown.markdown(text)
        # 简单提取文本（去掉 HTML 标签）
        import re
        text_content = re.sub(r'<[^>]+>', '', html)
        return [{"content": text_content, "metadata": {"source": filename, "type": ".md"}}]

    @staticmethod
    def parse_pdf(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 PDF 文件"""
        try:
            from pypdf import PdfReader
        except ImportError:
            return [{"content": "[请安装 pypdf 库]", "metadata": {"source": filename, "type": ".pdf"}}]

        results = []
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            reader = PdfReader(temp_path)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    results.append({
                        "content": text,
                        "metadata": {"source": filename, "type": ".pdf", "page": i + 1}
                    })
        finally:
            os.unlink(temp_path)

        return results if results else [{"content": "[PDF 无可提取文本]", "metadata": {"source": filename, "type": ".pdf"}}]

    @staticmethod
    def parse_html(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 HTML 文件"""
        from bs4 import BeautifulSoup
        text = file_content.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(text, 'html.parser')
        # 提取文本
        for script in soup(["script", "style"]):
            script.decompose()
        text_content = soup.get_text(separator='\n')
        return [{"content": text_content, "metadata": {"source": filename, "type": ".html"}}]

    @staticmethod
    def parse_docx(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 Word 文档"""
        try:
            from docx import Document
        except ImportError:
            return [{"content": "[请安装 python-docx 库]", "metadata": {"source": filename, "type": ".docx"}}]

        results = []
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            doc = Document(temp_path)
            # 按段落收集文本
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            if paragraphs:
                results.append({
                    "content": "\n".join(paragraphs),
                    "metadata": {"source": filename, "type": ".docx"}
                })
        finally:
            os.unlink(temp_path)

        return results if results else [{"content": "[DOCX 无可提取文本]", "metadata": {"source": filename, "type": ".docx"}}]

    @staticmethod
    def parse_excel(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 Excel 文件"""
        try:
            import pandas as pd
        except ImportError:
            return [{"content": "[请安装 pandas 库]", "metadata": {"source": filename, "type": ".xlsx"}}]

        results = []
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            # 读取所有 sheet
            excel_file = pd.ExcelFile(temp_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                # 转换为文本
                text_content = df.to_string(index=False)
                if text_content.strip():
                    results.append({
                        "content": f"Sheet: {sheet_name}\n{text_content}",
                        "metadata": {"source": filename, "type": ".xlsx", "sheet": sheet_name}
                    })
        finally:
            os.unlink(temp_path)

        return results if results else [{"content": "[Excel 无可提取数据]", "metadata": {"source": filename, "type": ".xlsx"}}]

    @staticmethod
    def parse_csv(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 CSV 文件"""
        try:
            import pandas as pd
        except ImportError:
            return [{"content": "[请安装 pandas 库]", "metadata": {"source": filename, "type": ".csv"}}]

        results = []
        with tempfile.NamedTemporaryFile(suffix='.csv', delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            df = pd.read_csv(temp_path)
            text_content = df.to_string(index=False)
            if text_content.strip():
                results.append({
                    "content": text_content,
                    "metadata": {"source": filename, "type": ".csv", "rows": len(df)}
                })
        finally:
            os.unlink(temp_path)

        return results if results else [{"content": "[CSV 无可提取数据]", "metadata": {"source": filename, "type": ".csv"}}]

    @staticmethod
    def parse_eml(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 EML 邮件文件"""
        try:
            import email
        except ImportError:
            return [{"content": "[Email 解析失败]", "metadata": {"source": filename, "type": ".eml"}}]

        results = []
        try:
            msg = email.message_from_bytes(file_content)
            # 提取主题
            subject = msg.get('subject', '')
            # 提取正文
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == 'text/plain':
                        body = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                        break
            else:
                body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')

            content = f"主题: {subject}\n\n{body}"
            results.append({
                "content": content,
                "metadata": {"source": filename, "type": ".eml", "subject": subject}
            })
        except Exception as e:
            results.append({
                "content": f"[EML 解析失败: {str(e)}]",
                "metadata": {"source": filename, "type": ".eml"}
            })

        return results

    @staticmethod
    def parse_msg(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 MSG 邮件文件（Outlook）"""
        # MSG 文件需要 outlook-msgfile 库，暂简单处理
        return [{"content": "[MSG 格式暂不支持，请转换为 EML]", "metadata": {"source": filename, "type": ".msg"}}]

    @staticmethod
    def parse_pptx(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 PPTX 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            return [{"content": "[请安装 python-pptx 库]", "metadata": {"source": filename, "type": ".pptx"}}]

        results = []
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            prs = Presentation(temp_path)
            slides_content = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                if texts:
                    slides_content.append(f"幻灯片 {i+1}:\n" + "\n".join(texts))

            if slides_content:
                results.append({
                    "content": "\n\n".join(slides_content),
                    "metadata": {"source": filename, "type": ".pptx", "slides": len(prs.slides)}
                })
        finally:
            os.unlink(temp_path)

        return results if results else [{"content": "[PPTX 无可提取文本]", "metadata": {"source": filename, "type": ".pptx"}}]

    @staticmethod
    def parse_xml(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 XML 文件"""
        try:
            import xml.etree.ElementTree as ET
        except ImportError:
            return [{"content": "[XML 解析失败]", "metadata": {"source": filename, "type": ".xml"}}]

        results = []
        try:
            text = file_content.decode('utf-8', errors='ignore')
            root = ET.fromstring(text)
            # 提取所有文本内容
            text_content = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    text_content.append(elem.text.strip())
            if text_content:
                results.append({
                    "content": "\n".join(text_content),
                    "metadata": {"source": filename, "type": ".xml"}
                })
        except Exception as e:
            results.append({
                "content": f"[XML 解析失败: {str(e)}]",
                "metadata": {"source": filename, "type": ".xml"}
            })

        return results if results else [{"content": "[XML 无可提取内容]", "metadata": {"source": filename, "type": ".xml"}}]

    @staticmethod
    def parse_epub(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析 EPUB 电子书"""
        # EPUB 解析需要 epub 库，暂返回提示
        return [{"content": "[EPUB 格式暂不支持，建议转换为 TXT]", "metadata": {"source": filename, "type": ".epub"}}]

    @staticmethod
    def parse_image(file_content: bytes, filename: str) -> List[Dict[str, Any]]:
        """解析图片文件 - 使用 OCR 提取文字"""
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            return [{"content": "[请安装 pytesseract 和 pillow 库]", "metadata": {"source": filename, "type": "image"}}]

        results = []
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(file_content)
            temp_path = f.name

        try:
            image = Image.open(temp_path)
            # OCR 提取文字
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            if text.strip():
                results.append({
                    "content": text,
                    "metadata": {"source": filename, "type": "image", "format": filename.split('.')[-1]}
                })
        finally:
            os.unlink(temp_path)

        return results if results else [{"content": "[图片无文字]", "metadata": {"source": filename, "type": "image"}}]
